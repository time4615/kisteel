import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_images_from_shape(shape, slide_no, image_count, output_dir):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image = shape.image
        image_bytes = image.blob
        image_filename = f"image_{slide_no + 1}_{image_count}.{image.ext}"
        image_path = os.path.join(output_dir, image_filename)
        
        with open(image_path, "wb") as f:
            f.write(image_bytes)
        print(f"Extracted: {image_filename}")
        return 1
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        count = 0
        for s in shape.shapes:
            count += extract_images_from_shape(s, slide_no, image_count + count, output_dir)
        return count
    return 0

def extract_all_images(pptx_path, output_dir):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    prs = Presentation(pptx_path)
    image_count = 0
    
    for slide_no, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            image_count += extract_images_from_shape(shape, slide_no, image_count, output_dir)

if __name__ == "__main__":
    pptx_file = "(주)케이아이스틸 회사소개서 2025ver.pptx"
    output_directory = "images_all"
    extract_all_images(pptx_file, output_directory)
    print("Done extracting all images.")
