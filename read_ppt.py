from pptx import Presentation

prs = Presentation("(주)케이아이스틸 회사소개서 2025ver.pptx")
with open("ppt_content_utf8.txt", "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"--- Slide {i+1} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                f.write(shape.text + "\n")
